import os
import re
from datetime import datetime
from collections import defaultdict
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_LEFT
import nltk

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt', quiet=True)

class LectureNoteGenerator:
    def __init__(self):
        self.file_path = None
        self.file_type = None
        
        # Legal topic keywords for detection
        self.topic_keywords = {
            'Offer and Acceptance': ['offer', 'acceptance', 'invitation to treat', 'postal rule', 'unilateral contract'],
            'Consideration': ['consideration', 'benefit', 'detriment', 'sufficient', 'adequate', 'past consideration'],
            'Promissory Estoppel': ['estoppel', 'promissory estoppel', 'shield not sword', 'equitable'],
            'Intention to Create Legal Relations': ['intention', 'legal relations', 'domestic', 'commercial'],
            'Certainty of Terms': ['certainty', 'vague', 'uncertain', 'agreement to agree'],
            'Capacity': ['capacity', 'minor', 'mental incapacity', 'intoxication'],
            'Privity of Contract': ['privity', 'third party', 'rights of third parties']
        }
    
    def process_file(self, file_path, progress_callback=None):
        """Main entry point for generating lecture notes"""
        self.file_path = file_path
        self.file_type = self._detect_file_type()
        
        if progress_callback:
            progress_callback("Loading file...", 0.1)
        
        # Extract raw content
        if self.file_type == 'powerpoint':
            raw_content = self._extract_powerpoint_content()
        else:
            raw_content = self._extract_text_content()
        
        if progress_callback:
            progress_callback("Analyzing topics...", 0.4)
        
        # Organize by topics
        topics = self._organize_by_topics(raw_content)
        
        if progress_callback:
            progress_callback("Formatting notes...", 0.7)
        
        # Format as study guide
        formatted_notes = self._format_as_study_guide(topics)
        
        if progress_callback:
            progress_callback("Complete!", 1.0)
        
        return {
            'notes': formatted_notes,
            'topics': topics
        }
    
    def _detect_file_type(self):
        """Detect if file is PowerPoint or text"""
        return 'powerpoint' if self.file_path.endswith('.pptx') else 'text'
    
    def _extract_powerpoint_content(self):
        """Extract all text from PowerPoint"""
        prs = Presentation(self.file_path)
        all_text = []
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    all_text.append(shape.text.strip())
            
            # Include speaker notes
            if slide.has_notes_slide:
                notes_slide = slide.notes_slide
                if notes_slide.notes_text_frame:
                    all_text.append(notes_slide.notes_text_frame.text.strip())
        
        return ' '.join(all_text)
    
    def _extract_text_content(self):
        """Extract and clean text content"""
        with open(self.file_path, 'r', encoding='utf-8') as f:
            text = f.read()
        
        # Clean transcript
        return self._clean_transcript(text)
    
    def _clean_transcript(self, text):
        """Aggressively clean conversational transcript"""
        # Remove filler words and phrases
        fillers = [
            r'\bum+\b', r'\buh+\b', r'\bah+\b', r'\ber+\b', r'\blike\b',
            r'\byou know\b', r'\bI mean\b', r'\bbasically\b', r'\bactually\b',
            r'\bliterally\b', r'\bkind of\b', r'\bsort of\b', r'\bI think\b',
            r'\bI guess\b', r'\bI believe\b', r'\bI feel like\b', r'\bso yeah\b',
            r'\bokay so\b', r'\balright so\b', r'\bwell\b', r'\bso\b(?=\s+\w)',
            r'\bjust\b', r'\breally\b', r'\bvery\b', r'\bquite\b',
            r'\blet me\b', r'\blet\'s\b', r'\bgoing to\b', r'\bgonna\b',
            r'\bOkay,?\s*', r'\bAlright,?\s*', r'\bNow,?\s*'
        ]
        
        for filler in fillers:
            text = re.sub(filler, '', text, flags=re.IGNORECASE)
        
        # Remove repeated words
        text = re.sub(r'\b(\w+)\s+\1\b', r'\1', text, flags=re.IGNORECASE)
        
        # Clean up whitespace
        text = re.sub(r'\s+', ' ', text)
        text = re.sub(r'\s+([.,!?])', r'\1', text)
        
        return text.strip()
    
    def _organize_by_topics(self, text):
        """Organize content by legal topics"""
        # Split into sentences
        sentences = re.split(r'[.!?]+', text)
        sentences = [s.strip() for s in sentences if len(s.strip()) > 20]
        
        # Initialize topic structure
        topics = {}
        for topic_name in self.topic_keywords.keys():
            topics[topic_name] = {
                'definitions': [],
                'rules': [],
                'cases': [],
                'examples': [],
                'exceptions': []
            }
        
        # Classify each sentence
        for sentence in sentences:
            # Detect which topic this sentence belongs to
            detected_topics = self._detect_topics(sentence)
            
            # Classify content type
            content_type = self._classify_content_type(sentence)
            
            # Add to appropriate topic(s)
            for topic in detected_topics:
                if topic in topics:
                    topics[topic][content_type].append(sentence)
        
        # Remove empty topics
        topics = {k: v for k, v in topics.items() if any(v.values())}
        
        return topics
    
    def _detect_topics(self, sentence):
        """Detect which topics a sentence relates to"""
        sentence_lower = sentence.lower()
        detected = []
        
        for topic, keywords in self.topic_keywords.items():
            if any(keyword in sentence_lower for keyword in keywords):
                detected.append(topic)
        
        # If no topic detected, assign to first topic (general)
        return detected if detected else [list(self.topic_keywords.keys())[0]]
    
    def _classify_content_type(self, sentence):
        """Classify sentence as definition, rule, case, example, or exception"""
        sentence_lower = sentence.lower()
        
        # Check for legal cases
        if re.search(r'[A-Z][a-z]+\s+v\.?\s+[A-Z][a-z]+', sentence):
            return 'cases'
        
        # Check for definitions
        if any(word in sentence_lower for word in ['is defined as', 'means', 'refers to', 'is the', 'are called']):
            return 'definitions'
        
        # Check for exceptions
        if any(word in sentence_lower for word in ['exception', 'unless', 'however', 'but', 'except']):
            return 'exceptions'
        
        # Check for examples
        if any(word in sentence_lower for word in ['for example', 'e.g.', 'such as', 'for instance']):
            return 'examples'
        
        # Check for rules (must, should, requirement, rule)
        if any(word in sentence_lower for word in ['must', 'should', 'requirement', 'rule', 'principle']):
            return 'rules'
        
        # Default to rules
        return 'rules'
    
    def _format_as_study_guide(self, topics):
        """Format organized topics as a study guide"""
        lines = []
        
        # Header
        lines.append("=" * 70)
        lines.append("LECTURE NOTES - STUDY GUIDE")
        lines.append("=" * 70)
        lines.append(f"Source: {os.path.basename(self.file_path)}")
        lines.append(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M')}")
        lines.append("")
        lines.append("Organized by legal topics for effective studying.")
        lines.append("")
        
        # Process each topic
        for topic_name, content in topics.items():
            lines.append("")
            lines.append("=" * 70)
            lines.append(f"TOPIC: {topic_name.upper()}")
            lines.append("=" * 70)
            lines.append("")
            
            # Definitions
            if content['definitions']:
                lines.append("DEFINITIONS:")
                lines.append("-" * 40)
                for definition in content['definitions'][:5]:  # Limit to 5
                    lines.append(f"• {definition}")
                lines.append("")
            
            # Legal Rules
            if content['rules']:
                lines.append("LEGAL RULES & PRINCIPLES:")
                lines.append("-" * 40)
                for rule in content['rules'][:10]:  # Limit to 10
                    lines.append(f"• {rule}")
                lines.append("")
            
            # Key Cases
            if content['cases']:
                lines.append("KEY CASES:")
                lines.append("-" * 40)
                for case in content['cases'][:8]:  # Limit to 8
                    # Format case name in bold
                    case_formatted = self._format_case_citation(case)
                    lines.append(f"• {case_formatted}")
                lines.append("")
            
            # Exceptions
            if content['exceptions']:
                lines.append("EXCEPTIONS & SPECIAL RULES:")
                lines.append("-" * 40)
                for exception in content['exceptions'][:5]:
                    lines.append(f"• {exception}")
                lines.append("")
            
            # Examples
            if content['examples']:
                lines.append("PRACTICAL EXAMPLES:")
                lines.append("-" * 40)
                for example in content['examples'][:5]:
                    lines.append(f"• {example}")
                lines.append("")
        
        # Study tips
        lines.append("")
        lines.append("=" * 70)
        lines.append("STUDY TIPS")
        lines.append("=" * 70)
        lines.append("• Review each topic's definitions first")
        lines.append("• Memorize key case names and principles")
        lines.append("• Understand exceptions to general rules")
        lines.append("• Practice applying rules to examples")
        lines.append("")
        
        return '\n'.join(lines)
    
    def _format_case_citation(self, text):
        """Format case citation with bold name"""
        # Extract case name
        match = re.search(r'([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)\s+v\.?\s+([A-Z][a-z]+(?:\s+[A-Z][a-z]+)*)', text)
        
        if match:
            case_name = f"{match.group(1)} v {match.group(2)}"
            
            # Extract year if present
            year_match = re.search(r'\((\d{4})\)', text)
            if year_match:
                case_name += f" ({year_match.group(1)})"
            
            # Replace in original text
            return text.replace(match.group(0), f"**{case_name}**")
        
        return text
    
    def export_to_pdf(self, data, filename):
        """Export lecture notes to PDF"""
        doc = SimpleDocTemplate(filename, pagesize=letter)
        styles = getSampleStyleSheet()
        story = []
        
        # Parse the notes text
        notes_text = data.get('notes', '')
        lines = notes_text.split('\n')
        
        for line in lines:
            if line.startswith('TOPIC:'):
                story.append(Spacer(1, 12))
                story.append(Paragraph(line, styles['Heading1']))
            elif line.startswith(('DEFINITIONS:', 'LEGAL RULES:', 'KEY CASES:', 'EXCEPTIONS:', 'EXAMPLES:')):
                story.append(Paragraph(f"<b>{line}</b>", styles['Heading2']))
            elif line.strip():
                story.append(Paragraph(line, styles['Normal']))
                story.append(Spacer(1, 4))
        
        doc.build(story)

# For backward compatibility
ConceptualAssistant = LectureNoteGenerator
